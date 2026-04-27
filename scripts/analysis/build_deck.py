"""Generate the Track-3 presentation deck as a .pptx file.

This mirrors submission/presentation_outline.md — seven 16:9 slides, ~30s each,
covering all five Rubric Section 3 questions (problem / importance / hardness /
prior work + differentiator / approach + results + limitations). Slide 7 embeds
submission/figures/per_split_scores.png as the Rubric Section 2 deliverable.

Usage:
    .venv/bin/python scripts/analysis/build_deck.py

Output:
    submission/presentation.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[2]
OUTLINE_PATH = REPO_ROOT / "submission" / "presentation_outline.md"
FIGURE_PATH = REPO_ROOT / "submission" / "figures" / "per_split_scores.png"
OUTPUT_PATH = REPO_ROOT / "submission" / "presentation.pptx"

# Princeton-ish palette with accessible contrast on white.
COLOR_TITLE = RGBColor(0x0F, 0x1B, 0x2D)        # deep navy
COLOR_BODY = RGBColor(0x1F, 0x2A, 0x3A)         # near-black
COLOR_ACCENT = RGBColor(0xE7, 0x72, 0x1E)       # Princeton-ish orange
COLOR_SUBTLE = RGBColor(0x55, 0x5F, 0x6D)       # muted slate
COLOR_RULE = RGBColor(0xE7, 0x72, 0x1E)
COLOR_TABLE_HEADER = RGBColor(0x0F, 0x1B, 0x2D)
COLOR_TABLE_BODY = RGBColor(0xF6, 0xF2, 0xEC)
COLOR_TABLE_HIGHLIGHT = RGBColor(0xFD, 0xEF, 0xDA)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _add_title(slide, text: str, subtitle: str | None = None) -> None:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = COLOR_TITLE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(16)
        r2.font.color.rgb = COLOR_SUBTLE

    # Thin accent rule under the title block.
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.35), Inches(2.25), Emu(38100)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = COLOR_RULE
    rule.line.fill.background()


def _add_footer(slide, slide_num: int, total: int, label: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.333), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r1 = p.add_run()
    r1.text = f"CityLearn 2023 Track-3  ·  {label}"
    r1.font.size = Pt(10)
    r1.font.color.rgb = COLOR_SUBTLE
    r2 = p.add_run()
    r2.text = f"     {slide_num} / {total}"
    r2.font.size = Pt(10)
    r2.font.color.rgb = COLOR_SUBTLE


def _add_bullets(
    slide,
    items: list[str | tuple[str, str]],
    *,
    left: float = 0.5,
    top: float = 1.8,
    width: float = 12.333,
    height: float = 5.0,
    body_size: int = 18,
    lead_size: int | None = None,
) -> None:
    """Render a bullet list. Items can be plain strings or (lead, trailer)
    tuples where the lead is rendered bold.
    """
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        if isinstance(item, tuple):
            lead, trailer = item
            r_bullet = p.add_run()
            r_bullet.text = "• "
            r_bullet.font.size = Pt(body_size)
            r_bullet.font.color.rgb = COLOR_ACCENT
            r_bullet.font.bold = True
            r_lead = p.add_run()
            r_lead.text = lead
            r_lead.font.size = Pt(lead_size or body_size)
            r_lead.font.bold = True
            r_lead.font.color.rgb = COLOR_BODY
            if trailer:
                r_trail = p.add_run()
                r_trail.text = trailer
                r_trail.font.size = Pt(body_size)
                r_trail.font.color.rgb = COLOR_BODY
        else:
            r_bullet = p.add_run()
            r_bullet.text = "• "
            r_bullet.font.size = Pt(body_size)
            r_bullet.font.color.rgb = COLOR_ACCENT
            r_bullet.font.bold = True
            r_text = p.add_run()
            r_text.text = item
            r_text.font.size = Pt(body_size)
            r_text.font.color.rgb = COLOR_BODY
        p.space_after = Pt(10)


def _set_notes(slide, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text.strip()


def _add_reward_table(slide) -> None:
    rows, cols = 4, 3
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(12.333)
    height = Inches(4.2)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.columns[0].width = Inches(2.1)
    table.columns[1].width = Inches(6.5)
    table.columns[2].width = Inches(3.733)

    headers = ["Variant", "What it is", "Why"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TABLE_HEADER
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = h
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

    data = [
        (
            "reward_v0\n(baseline)",
            "CityLearn's built-in default (negative net-consumption proxy).",
            "Sanity baseline; the training signal does not match the competition scoring rule.",
            False,
        ),
        (
            "reward_v1",
            (
                "8-term weighted penalty mirroring the official 2023 "
                "scoring weights:\n"
                "0.30·comfort + 0.15·outage_comfort + 0.15·outage_unserved + "
                "0.10·carbon + 0.075·{ramping, load_factor, daily_peak, all_time_peak}"
            ),
            "Train on the same thing the competition scores on.",
            True,
        ),
        (
            "reward_v2",
            "reward_v1 + 0.05·|ΔSoC| + 0.05·1[sign_flip]",
            "Battery smoothness penalties — discourages bang-bang charge/discharge; "
            "empirically improves cross-split generalization.",
            True,
        ),
    ]
    for r_idx, (variant, what, why, bold_variant) in enumerate(data, start=1):
        highlight = bold_variant
        for c_idx, text in enumerate([variant, what, why]):
            cell = table.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_TABLE_HIGHLIGHT if highlight else COLOR_TABLE_BODY
            tf = cell.text_frame
            tf.word_wrap = True
            # Replace newlines with explicit paragraphs for proper rendering.
            lines = text.split("\n")
            for i, line in enumerate(lines):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line
                run.font.size = Pt(12)
                run.font.color.rgb = COLOR_BODY
                if c_idx == 0:
                    run.font.bold = bold_variant
                    run.font.size = Pt(13)


def _add_results_numbers_block(slide) -> None:
    left = Inches(7.9)
    top = Inches(1.85)
    tb = slide.shapes.add_textbox(left, top, Inches(5.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True

    blocks = [
        ("phase_2 (3 buildings, held-out)", [
            ("RBC", "1.087"),
            ("PPO central", "0.873"),
            ("PPO DTDE", "0.793"),
            ("SAC central reward_v2", "0.653", True),
            ("SAC DTDE", "0.677"),
            ("CHESCA public ref", "0.562", False, True),
        ]),
        ("phase_3 (6 buildings, held-out — portable only)", [
            ("RBC", "1.114"),
            ("PPO DTDE", "0.843"),
            ("SAC DTDE", "0.774", True),
        ]),
    ]
    first = True
    for header, rows in blocks:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = header
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = COLOR_ACCENT
        p.space_after = Pt(4)

        for row in rows:
            label = row[0]
            value = row[1]
            best = row[2] if len(row) > 2 else False
            is_ref = row[3] if len(row) > 3 else False
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r_lbl = p2.add_run()
            r_lbl.text = f"  {label}: "
            r_lbl.font.size = Pt(13)
            r_lbl.font.color.rgb = COLOR_SUBTLE if is_ref else COLOR_BODY
            r_val = p2.add_run()
            r_val.text = value
            r_val.font.size = Pt(13)
            r_val.font.bold = True
            if best:
                r_val.font.color.rgb = COLOR_ACCENT
            elif is_ref:
                r_val.font.color.rgb = COLOR_SUBTLE
            else:
                r_val.font.color.rgb = COLOR_BODY
            if best:
                r_tag = p2.add_run()
                r_tag.text = "  (best)"
                r_tag.font.size = Pt(11)
                r_tag.font.italic = True
                r_tag.font.color.rgb = COLOR_ACCENT


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank = prs.slide_layouts[6]  # blank layout; we draw everything manually
    total_slides = 7

    # ---- Slide 1: Title + problem ----
    s1 = prs.slides.add_slide(blank)
    # Hero title (centered block)
    tb = s1.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(11.933), Inches(2.0))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "Generalizable RL for Neighborhood Battery Control"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = COLOR_TITLE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "CityLearn Challenge 2023 · Track 3 (RL Competition)"
    r2.font.size = Pt(20)
    r2.font.color.rgb = COLOR_ACCENT
    r2.font.bold = True

    # Accent rule
    rule = s1.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.45), Inches(2.25), Emu(38100)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = COLOR_ACCENT
    rule.line.fill.background()

    # Team
    tb2 = s1.shapes.add_textbox(Inches(0.7), Inches(4.7), Inches(11.933), Inches(1.5))
    tf2 = tb2.text_frame
    p3 = tf2.paragraphs[0]
    r3 = p3.add_run()
    r3.text = "Erik Dyer  ·  Ammaar Alam  ·  Grace Sun"
    r3.font.size = Pt(20)
    r3.font.color.rgb = COLOR_BODY
    p4 = tf2.add_paragraph()
    r4 = p4.add_run()
    r4.text = "COS 435 / ECE 433 — Spring 2026"
    r4.font.size = Pt(16)
    r4.font.color.rgb = COLOR_SUBTLE

    # Problem statement block
    tb3 = s1.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(11.933), Inches(1.3))
    tf3 = tb3.text_frame
    tf3.word_wrap = True
    p5 = tf3.paragraphs[0]
    r5 = p5.add_run()
    r5.text = (
        "Problem: control a shared neighborhood battery across 3 buildings with "
        "solar PV and stochastic demand. Each hour the policy emits a per-building "
        "continuous action in [−1, +1] to jointly reduce cost, carbon, peak, "
        "discomfort, and outage risk."
    )
    r5.font.size = Pt(15)
    r5.font.color.rgb = COLOR_BODY
    r5.font.italic = True

    _add_footer(s1, 1, total_slides, "Title + problem")
    _set_notes(
        s1,
        "We control a neighborhood-scale battery storage system across three "
        "buildings with solar generation and stochastic electricity demand. At "
        "each hour, the agent outputs a continuous action in [-1, +1] per building "
        "— negative discharges, positive charges — to optimize a composite score "
        "across cost, carbon, peak demand, discomfort, and thermal resilience.",
    )

    # ---- Slide 2: Why it matters ----
    s2 = prs.slides.add_slide(blank)
    _add_title(s2, "Why it matters", "Buildings are a first-order decarbonization lever")
    _add_bullets(
        s2,
        [
            ("Buildings ≈ 28% of global CO₂ emissions", "  (IEA)"),
            (
                "Distributed storage + PV",
                " is a key grid-decarbonization lever — shifts peak, absorbs renewables.",
            ),
            (
                "CityLearn Challenge 2023",
                " — realistic, standardized benchmark for multi-building coordination.",
            ),
            (
                "Better battery coordination",
                " → lower peak loads, cheaper bills, less carbon. Real-world stakes.",
            ),
        ],
        top=1.9,
        body_size=20,
    )
    _add_footer(s2, 2, total_slides, "Why it matters")
    _set_notes(
        s2,
        "Better battery coordination directly translates to lower peak loads, "
        "cheaper bills, and less carbon — real-world stakes. The CityLearn "
        "benchmark gives us a controlled testbed to study RL for a problem where "
        "rule-based and MPC methods currently dominate.",
    )

    # ---- Slide 3: Why it's hard ----
    s3 = prs.slides.add_slide(blank)
    _add_title(s3, "Why it is hard", "Naive approaches fail for concrete, measurable reasons")
    _add_bullets(
        s3,
        [
            ("Non-stationary demand + weather", " — distribution shifts across seasons and hours."),
            ("Partial observability", " — no perfect load forecast, imperfect thermal state."),
            ("Multi-building coordination, long horizons", " — credit assignment is brutal."),
            (
                "Reward specification is not given",
                " — competition ships an evaluation rule, not a training signal.",
            ),
            (
                "Naive RBC = 1.02 on public_dev",
                "  — strong non-RL baseline. Exploitable structure, little adaptivity.",
            ),
        ],
        top=1.9,
        body_size=19,
    )
    _add_footer(s3, 3, total_slides, "Why it is hard")
    _set_notes(
        s3,
        "Rule-based control does a reasonable job because the domain has "
        "exploitable structure, but it cannot adapt to unusual conditions. RL "
        "promises better generalization, but training is unstable across seeds, "
        "reward misspecification dominates early results, and the long-horizon "
        "credit assignment is brutal. A naive single-seed PPO run looks better "
        "than RBC but regresses on held-out splits.",
    )

    # ---- Slide 4: Prior work + our differentiator ----
    s4 = prs.slides.add_slide(blank)
    _add_title(
        s4,
        "Prior work  +  our differentiator",
        "Track 3 is already solved if the metric is leaderboard rank",
    )
    _add_bullets(
        s4,
        [
            (
                "CityLearn 2023 winner: CHESCA",
                " — 0.562 public / 0.565 private (hierarchical optimization, not RL).",
            ),
            (
                "Most leaderboard entries",
                " are rule-based or MPC variants — RL is under-represented.",
            ),
            ("Our contribution: a systematic ablation, not a new algorithm.", ""),
            ("2 algorithms  ×  2 control topologies  ×  3 reward variants", ""),
            ("Multi-seed runs (n = 3–10) with 95% CI error bars", ""),
            ("Cross-split generalization", ": public_dev → phase_2 (3b) → phase_3 (6b held-out)"),
        ],
        top=1.85,
        body_size=18,
    )
    _add_footer(s4, 4, total_slides, "Prior work + differentiator")
    _set_notes(
        s4,
        "Track 3 has already been solved if the metric is leaderboard rank — "
        "CHESCA won with hierarchical optimization. Nobody has done a clean "
        "comparison of RL architectures and rewards with seeded error bars, "
        "though. That's the gap we fill. We don't claim to beat CHESCA; we claim "
        "to map which design choices matter for RL practitioners who come next.",
    )

    # ---- Slide 5: Approach - architecture + algorithms ----
    s5 = prs.slides.add_slide(blank)
    _add_title(
        s5,
        "Approach — architecture + algorithms",
        "Two control topologies, two RL algorithms",
    )
    _add_bullets(
        s5,
        [
            (
                "Centralized",
                ": one controller observes all 3 buildings, emits a 3-vector action. "
                "Wired for exactly 3 buildings.",
            ),
            (
                "Shared-DTDE",
                ": one policy network, per-building observation windows, decentralized "
                "rollout with shared parameters — size-agnostic.",
            ),
            (
                "PPO",
                " — on-policy, clipped surrogate objective. From-scratch for DTDE; "
                "Stable-Baselines-3 for centralized.",
            ),
            ("SAC", " — off-policy, entropy-regularized. Custom for both topologies."),
            (
                "Only shared-DTDE transfers to the 6-building phase_3 held-out cluster",
                " — a centralized policy's input/output shapes are fundamentally non-portable.",
            ),
        ],
        top=1.85,
        body_size=17,
    )
    _add_footer(s5, 5, total_slides, "Approach — architecture")
    _set_notes(
        s5,
        "We implemented PPO from scratch on a shared policy net with per-building "
        "embeddings for the DTDE variants, and used Stable-Baselines-3 for "
        "centralized PPO. SAC is custom for both topologies. The centralized "
        "architecture gets full joint-action credit assignment; the DTDE version "
        "is the only one that transfers to the 6-building phase-3 held-out split "
        "— centralized policies are fundamentally non-portable there.",
    )

    # ---- Slide 6: Approach - reward design ----
    s6 = prs.slides.add_slide(blank)
    _add_title(s6, "Approach — reward design", "Train on the same thing the competition scores on")
    _add_reward_table(s6)

    # Below-table footnote
    tb = s6.shapes.add_textbox(Inches(0.5), Inches(6.25), Inches(12.333), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = (
        "Reward-variant ablation axis run on SAC-central only (compute-bound). "
        "SoC = battery state of charge; sign_flip indicator penalizes charge→discharge "
        "reversal between consecutive steps."
    )
    r.font.size = Pt(11)
    r.font.italic = True
    r.font.color.rgb = COLOR_SUBTLE

    _add_footer(s6, 6, total_slides, "Reward design")
    _set_notes(
        s6,
        "Instead of inventing weights, reward_v1 literally mirrors the 8-term "
        "weighted penalty the competition uses to score submissions, so the "
        "agent is trained directly on the objective it's evaluated on. reward_v2 "
        "adds two battery-behavior terms that penalize SoC swings and charge-"
        "discharge sign flips between consecutive steps. Without them, SAC learns "
        "bang-bang policies that are locally optimal on public_dev but brittle "
        "on held-out splits; with them, the policy is steadier and generalizes "
        "better. The reward-ablation axis is run only on SAC-central — we didn't "
        "have compute to matched-ablate PPO.",
    )

    # ---- Slide 7: Results + limitations ----
    s7 = prs.slides.add_slide(blank)
    _add_title(
        s7,
        "Results  +  limitations",
        "Held-out average_score (↓ better) across 6 released evaluation datasets",
    )
    # Figure on left (wider)
    if FIGURE_PATH.exists():
        s7.shapes.add_picture(
            str(FIGURE_PATH),
            Inches(0.5),
            Inches(1.85),
            width=Inches(7.2),
            height=Inches(4.5),
        )
    else:
        tb = s7.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(7.0), Inches(4.5))
        tb.text_frame.text = f"[figure missing: {FIGURE_PATH}]"

    # Numbers on right
    _add_results_numbers_block(s7)

    # Takeaway strip at bottom
    tb = s7.shapes.add_textbox(Inches(0.5), Inches(6.45), Inches(12.333), Inches(0.6))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r_lbl = p.add_run()
    r_lbl.text = "Takeaway: "
    r_lbl.font.size = Pt(14)
    r_lbl.font.bold = True
    r_lbl.font.color.rgb = COLOR_ACCENT
    r_text = p.add_run()
    r_text.text = (
        "shared-DTDE is the only architecture that transfers to a different "
        "building count; centralized SAC reward_v2 is the best within-size variant. "
        "CHESCA's 0.562 is benchmark context (different eval server), not a head-to-head number."
    )
    r_text.font.size = Pt(13)
    r_text.font.color.rgb = COLOR_BODY

    _add_footer(s7, 7, total_slides, "Results + limitations")
    _set_notes(
        s7,
        "The headline panel is the held-out phase_2_online_eval split — same "
        "3-building setup as training but weather and demand the agent has never "
        "seen. RBC sits at 1.09, PPO central at 0.87, PPO-DTDE at 0.79, and SAC "
        "variants cluster at 0.65-0.68. The dashed line at 0.562 is CHESCA's "
        "2023 public-leaderboard score — benchmark context, not a head-to-head "
        "number, because CHESCA was scored on the original competition server "
        "we cannot re-run.\n\n"
        "The harder test is phase_3 — held-out AND a different 6-building "
        "cluster, not the 3 we trained on. Only the shared-DTDE variants can "
        "execute here: a centralized policy's input and output layers are wired "
        "for exactly 3 buildings, so feeding it a 6-building observation is a "
        "shape mismatch, not a degradation. SAC-DTDE generalizes best at 0.774, "
        "PPO-DTDE at 0.843. This is our cleanest architectural claim: if "
        "deployment might see a different building count than training, choose "
        "shared-DTDE.\n\n"
        "Limitations: single public_dev tuning split was used to select reward "
        "variants (overfitting risk); compute-bound so we capped PPO-DTDE and "
        "SAC-DTDE at 10 and 3 seeds respectively; reward-variant ablation was "
        "run only on SAC-central, not on PPO-central. CHESCA's line is "
        "reference, not a competitive bar we're claiming to have cleared.",
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    path = build_presentation()
    print(f"wrote: {path.relative_to(REPO_ROOT)}  ({path.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
